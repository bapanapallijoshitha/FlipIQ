from flask import Flask, render_template, request, jsonify, session
import json, re, io, base64, os, uuid, datetime
import PyPDF2
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract
from groq import Groq

app = Flask(__name__)
app.secret_key = "flipiq_secret_key_2024"
GROQ_API_KEY = ""  # Add your key here

if not GROQ_API_KEY:
    print("\n" + "="*60)
    print("  ERROR: Groq API key is missing!")
    print("  1. Go to console.groq.com and sign up free")
    print("  2. Click API Keys → Create API Key")
    print("  3. Paste it in flipiq.py line 16")
    print("="*60 + "\n")

client = Groq(api_key=GROQ_API_KEY)

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

quiz_history = []


# ══════════════════════════════════════════════════════════════════════════════
# FILE EXTRACTION FUNCTIONS
# ══════════════════════════════════════════════════════════════════════════════

def extract_pdf_text(file_bytes: bytes) -> str:
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    return "\n".join(p.extract_text() or "" for p in reader.pages).strip()


def extract_pptx_text(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines).strip()


def extract_docx_text(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()


def extract_image_text(file_bytes: bytes) -> str:
    image = Image.open(io.BytesIO(file_bytes))
    return pytesseract.image_to_string(image).strip()


def extract_text_by_type(file_bytes: bytes, filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext == "pdf":
        return extract_pdf_text(file_bytes)
    elif ext in ("ppt", "pptx"):
        return extract_pptx_text(file_bytes)
    elif ext == "docx":
        return extract_docx_text(file_bytes)
    elif ext in ("jpg", "jpeg", "png", "bmp", "tiff", "webp"):
        return extract_image_text(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")


# ══════════════════════════════════════════════════════════════════════════════
# AUTO-DETECT DIFFICULTY
# ══════════════════════════════════════════════════════════════════════════════

def auto_detect_difficulty(notes: str) -> str:
    """Ask Groq to assess the complexity of the material."""
    try:
        response = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[{
                "role": "user",
                "content": f"""Analyze this study material and respond with ONLY one word:
- "easy" if it contains basic facts, simple definitions, introductory concepts
- "medium" if it contains moderate concepts requiring understanding
- "hard" if it contains advanced theory, complex analysis, or technical depth

Study material (first 2000 chars):
\"\"\"{notes[:2000]}\"\"\"

Respond with only: easy, medium, or hard"""
            }],
            max_tokens=10,
            temperature=0.1,
        )
        result = response.choices[0].message.content.strip().lower()
        if result in ("easy", "medium", "hard"):
            return result
        return "medium"
    except:
        return "medium"


# ══════════════════════════════════════════════════════════════════════════════
# QUIZ GENERATION
# ══════════════════════════════════════════════════════════════════════════════

def generate_quiz(notes: str, num_q: int, difficulty: str, focus: str) -> dict:
    focus_map = {
        "general":     "covering a broad range of topics",
        "definitions": "focusing on definitions and key terms",
        "concepts":    "focusing on core concepts and principles",
        "examples":    "focusing on examples and real-world applications",
    }
    diff_map = {
        "easy":   "basic recall questions suitable for beginners",
        "medium": "comprehension questions requiring understanding of concepts",
        "hard":   "application and analysis questions requiring critical thinking",
    }

    prompt = f"""You are an expert professor and educator. Based on the study material below,
generate exactly {num_q} multiple choice questions {focus_map.get(focus, 'covering a broad range of topics')}.
Difficulty: {diff_map.get(difficulty, 'comprehension questions')}

Study material:
\"\"\"
{notes[:8000]}
\"\"\"

Return ONLY a valid JSON object — no markdown, no explanation, no extra text before or after:
{{
  "subject": "short topic name (2-4 words)",
  "questions": [
    {{
      "id": 1,
      "question": "Question text?",
      "options": ["A) ...", "B) ...", "C) ...", "D) ..."],
      "correct": "A",
      "explanation": "One or two sentences explaining why this is correct.",
      "real_world_example": "A concrete real-world example or analogy that helps understand this concept easily."
    }}
  ]
}}

Rules:
- Exactly 4 options per question labeled A) B) C) D)
- "correct" must be only a single letter: A, B, C, or D
- Base every question strictly on the provided material
- "real_world_example" must be a simple, relatable everyday example a student can instantly understand
- Generate exactly {num_q} questions
- Output ONLY the JSON object, nothing else"""

    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=3000,
    )

    raw = response.choices[0].message.content
    raw = re.sub(r"```(?:json)?", "", raw).strip()
    match = re.search(r'\{.*\}', raw, re.DOTALL)
    if not match:
        raise ValueError("Could not parse AI response. Please try again.")
    return json.loads(match.group())


# ══════════════════════════════════════════════════════════════════════════════
# ROUTES
# ══════════════════════════════════════════════════════════════════════════════

@app.route("/")
def home():
    return render_template("index.html")


@app.route("/detect-difficulty", methods=["POST"])
def detect_difficulty():
    """Auto-detect difficulty from uploaded content."""
    try:
        notes = ""
        if "file" in request.files and request.files["file"].filename:
            f = request.files["file"]
            notes = extract_text_by_type(f.read(), f.filename)
        if not notes.strip():
            notes = request.form.get("notes_text", "")
        if len(notes.strip()) < 50:
            return jsonify({"difficulty": "medium"})
        difficulty = auto_detect_difficulty(notes)
        return jsonify({"difficulty": difficulty})
    except Exception as e:
        return jsonify({"difficulty": "medium", "error": str(e)})


@app.route("/generate", methods=["POST"])
def generate():
    try:
        num_q = int(request.form.get("num_questions", 5))
        diff  = request.form.get("difficulty", "medium")
        focus = request.form.get("focus", "general")
        notes = ""
        source_name = "Pasted text"

        # Extract from uploaded file
        if "file" in request.files and request.files["file"].filename:
            f = request.files["file"]
            source_name = f.filename
            file_bytes = f.read()
            notes = extract_text_by_type(file_bytes, f.filename)
            if not notes.strip():
                return jsonify({"error": "Could not extract text from file. Please try a different file."}), 400

        # Fallback to pasted text
        if not notes.strip():
            notes = request.form.get("notes_text", "")

        if len(notes.strip()) < 50:
            return jsonify({"error": "Please provide more study material (at least a few sentences)."}), 400

        quiz = generate_quiz(notes, num_q, diff, focus)
        quiz["difficulty"] = diff
        quiz["source"]     = source_name
        quiz["id"]         = str(uuid.uuid4())[:8]
        quiz["timestamp"]  = datetime.datetime.now().strftime("%d %b %Y, %I:%M %p")
        quiz["num_q"]      = num_q

        # Save to history (keep last 10)
        quiz_history.insert(0, {
            "id":        quiz["id"],
            "subject":   quiz["subject"],
            "source":    source_name,
            "difficulty": diff,
            "num_q":     num_q,
            "timestamp": quiz["timestamp"],
            "quiz":      quiz,
        })
        if len(quiz_history) > 10:
            quiz_history.pop()

        return jsonify(quiz)

    except json.JSONDecodeError:
        return jsonify({"error": "AI returned unexpected format. Please try again."}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/history")
def history():
    return jsonify(quiz_history)


@app.route("/history/<quiz_id>")
def get_history_quiz(quiz_id):
    for item in quiz_history:
        if item["id"] == quiz_id:
            return jsonify(item["quiz"])
    return jsonify({"error": "Quiz not found"}), 404


if __name__ == "__main__":
    app.run(debug=True)
